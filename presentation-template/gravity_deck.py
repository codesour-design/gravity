#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
GRAVITY DECK — generatore di slide aziendali con il tema Gravity.

Tema PowerPoint riutilizzabile per CodeSour / Gravity:
  - Cover e divisori di sezione su sfondo dark (#0A0A0A)
  - Slide di contenuto su sfondo light
  - Titoli in Oswald, testo in Inter (font incorporati nel file finale)
  - Colori brand: primary #3E00FB, secondary #FF4A1C

USO RAPIDO
----------
    from gravity_deck import GravityDeck

    d = GravityDeck()
    d.cover("Titolo del documento", "Sottotitolo / claim",
            meta="CodeSour · Giugno 2026")
    d.index([("Contesto e obiettivi"), ("Soluzione"), ("Valore")])
    d.section(1, "Contesto e obiettivi")
    d.content("Il problema", body="Testo introduttivo...",
              bullets=["Primo punto", "Secondo punto"], kicker="Contesto")
    d.cards("Tre pilastri", [
        {"title": "Coerenza", "text": "Una sola fonte di verità."},
        {"title": "Tracciabilità", "text": "Ogni scelta documentata."},
        {"title": "Velocità", "text": "Meno cicli di rilavorazione."},
    ])
    d.closing("buon lavoro", "Gravity · CodeSour")
    d.save("out.pptx")            # i font vengono incorporati automaticamente

Lo stesso file può essere guidato da un dizionario/JSON (vedi build_from_spec).
"""

import os
import shutil
import zipfile
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

# ─────────────────────────────────────────────────────────────────────────────
#  TOKEN DI BRAND
# ─────────────────────────────────────────────────────────────────────────────
PURPLE = RGBColor(0x3E, 0x00, 0xFB)   # primary  – Gravity blue elettrico
ORANGE = RGBColor(0xFF, 0x4A, 0x1C)   # secondary – Gravity orange
BLACK  = RGBColor(0x0A, 0x0A, 0x0A)   # sfondo dark
INK    = RGBColor(0x14, 0x14, 0x16)   # testo scuro su light
GREY   = RGBColor(0x6B, 0x6B, 0x72)   # testo secondario
FAINT  = RGBColor(0xB6, 0xB6, 0xBE)   # testo terziario / numeri pagina
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
PAPER  = RGBColor(0xFF, 0xFF, 0xFF)   # sfondo light
MIST   = RGBColor(0xF4, 0xF3, 0xFA)   # tint card (lavanda chiarissima)
LAV    = RGBColor(0xEC, 0xE9, 0xFF)   # lavanda
DARK2  = RGBColor(0x15, 0x15, 0x18)   # card su dark

TITLE_FONT = "Oswald"
BODY_FONT  = "Inter"

# 16:9
EMU_IN = 914400
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

ASSETS = os.path.join(os.path.dirname(os.path.abspath(__file__)), "assets")
FONTS  = os.path.join(os.path.dirname(os.path.abspath(__file__)), "fonts")


def _asset(name):
    return os.path.join(ASSETS, name)


def _hex(h):
    """'#RRGGBB' -> RGBColor."""
    h = h.lstrip("#")
    return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))


# ─────────────────────────────────────────────────────────────────────────────
#  HELPER DI BASSO LIVELLO
# ─────────────────────────────────────────────────────────────────────────────
def _solid(shape, color):
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()


def _bg(slide, color):
    r = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, SLIDE_H)
    _solid(r, color)
    r.shadow.inherit = False
    # manda lo sfondo dietro a tutto
    sp = r._element
    sp.getparent().remove(sp)
    slide.shapes._spTree.insert(2, sp)
    return r


def _txt(slide, x, y, w, h, anchor=MSO_ANCHOR.TOP, align=PP_ALIGN.LEFT, wrap=True):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = wrap
    tf.vertical_anchor = anchor
    for m in ("margin_left", "margin_right", "margin_top", "margin_bottom"):
        setattr(tf, m, 0)
    p = tf.paragraphs[0]
    p.alignment = align
    return tb, tf


def _run(p, text, font=BODY_FONT, size=14, color=INK, bold=False,
         italic=False, spacing=None):
    r = p.add_run()
    r.text = text
    r.font.name = font
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.italic = italic
    r.font.color.rgb = color
    if spacing is not None:                       # tracking in punti
        rPr = r._r.get_or_add_rPr()
        rPr.set("spc", str(int(spacing * 100)))
    return r


def _set_line_spacing(p, mult):
    p.line_spacing = mult


def _bullet_para(tf, text, first=False, size=15, color=INK,
                 marker=ORANGE, gap_before=8):
    p = tf.paragraphs[0] if first else tf.add_paragraph()
    p.space_after = Pt(gap_before)
    p.line_spacing = 1.18
    pPr = p._pPr if p._pPr is not None else p.get_or_add_pPr()
    pPr.set("marL", str(int(0.30 * EMU_IN)))
    pPr.set("indent", str(int(-0.30 * EMU_IN)))
    # bullet quadratino arancione (glyph preso da Arial per sicurezza)
    buClr = pPr.makeelement(qn("a:buClr"), {})
    srgb = pPr.makeelement(qn("a:srgbClr"), {"val": "%02X%02X%02X" % (marker[0], marker[1], marker[2])})
    buClr.append(srgb)
    buSzPct = pPr.makeelement(qn("a:buSzPct"), {"val": "70000"})
    buFont = pPr.makeelement(qn("a:buFont"), {"typeface": "Arial"})
    buChar = pPr.makeelement(qn("a:buChar"), {"char": "▪"})
    for el in (buClr, buSzPct, buFont, buChar):
        pPr.append(el)
    _run(p, text, font=BODY_FONT, size=size, color=color)
    return p


def _pic(slide, path, x, y, w=None, h=None):
    return slide.shapes.add_picture(path, x, y, width=w, height=h)


# colore di un RGBColor come tupla
def _rgb_tuple(c):
    return (c[0], c[1], c[2])


# ─────────────────────────────────────────────────────────────────────────────
#  DECK
# ─────────────────────────────────────────────────────────────────────────────
class GravityDeck:
    def __init__(self):
        self.prs = Presentation()
        self.prs.slide_width = SLIDE_W
        self.prs.slide_height = SLIDE_H
        self._blank = self.prs.slide_layouts[6]
        self._page = 0

    # — utilità interne —
    def _new(self):
        return self.prs.slides.add_slide(self._blank)

    def _footer(self, slide, dark=False):
        """Wordmark in basso a sx + numero pagina in basso a dx."""
        self._page += 1
        logo = "type_white.png" if dark else "type_purple.png"
        # wordmark ~ 1.05in di larghezza, ratio 357:120
        w = Inches(1.05)
        h = Emu(int(w * 120 / 357))
        _pic(slide, _asset(logo), Inches(0.6), Inches(7.5) - Inches(0.62), w=w, h=h)
        tb, tf = _txt(slide, Inches(11.6), Inches(6.86), Inches(1.13), Inches(0.4),
                      align=PP_ALIGN.RIGHT)
        _run(tf.paragraphs[0], "%02d" % self._page, font=BODY_FONT, size=10,
             color=FAINT if not dark else RGBColor(0x6A, 0x6A, 0x72))

    def _kicker(self, slide, text, x, y, w, color=ORANGE):
        if not text:
            return
        tb, tf = _txt(slide, x, y, w, Inches(0.34))
        _run(tf.paragraphs[0], text.upper(), font=TITLE_FONT, size=13,
             color=color, bold=True, spacing=2.5)

    # ── COVER (dark) ─────────────────────────────────────────────────────────
    def cover(self, title, subtitle=None, meta=None, kicker=None, stats=None):
        s = self._new()
        _bg(s, PURPLE)
        # monogramma grande sfumato a destra (viola scuro su viola)
        mk = _pic(s, _asset("mark_deep.png"), Inches(8.7), Inches(-1.1),
                  h=Inches(9.6))
        # monogramma piccolo in alto a sinistra (wordmark)
        w = Inches(1.7); h = Emu(int(w * 120 / 357))
        _pic(s, _asset("type_white.png"), Inches(0.85), Inches(0.8), w=w, h=h)
        # kicker
        if kicker:
            self._kicker(s, kicker, Inches(0.9), Inches(2.55), Inches(8),
                         color=ORANGE)
        # titolo
        tb, tf = _txt(s, Inches(0.85), Inches(2.85), Inches(9.4), Inches(2.8),
                      anchor=MSO_ANCHOR.TOP)
        p = tf.paragraphs[0]; p.line_spacing = 0.98
        _run(p, title, font=TITLE_FONT, size=58, color=WHITE, bold=True)
        # sottotitolo (più in alto se ci sono i numeri)
        sub_y = 4.4 if stats else 5.35
        if subtitle:
            tb2, tf2 = _txt(s, Inches(0.9), Inches(sub_y), Inches(8.8), Inches(1.0))
            p2 = tf2.paragraphs[0]; p2.line_spacing = 1.25
            _run(p2, subtitle, font=BODY_FONT, size=18, color=RGBColor(0xDA, 0xD3, 0xFF))
        # riga di numeri (come nella cover sorgente)
        if stats:
            n = len(stats); colw = 11.0 / n
            for i, (big, label) in enumerate(stats):
                x = 0.9 + i * colw
                tbb, tfb = _txt(s, Inches(x), Inches(5.35), Inches(colw - 0.2), Inches(0.85))
                _run(tfb.paragraphs[0], big, font=TITLE_FONT, size=40,
                     color=WHITE, bold=True)
                tbl, tfl = _txt(s, Inches(x), Inches(6.12), Inches(colw - 0.2), Inches(0.4))
                _run(tfl.paragraphs[0], label, font=BODY_FONT, size=12,
                     color=RGBColor(0xC6, 0xBF, 0xF2))
        # meta + dot arancione (solo se c'è il meta)
        if meta:
            dot = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(0.92), Inches(6.66),
                                     Inches(0.16), Inches(0.16))
            _solid(dot, ORANGE)
            tbm, tfm = _txt(s, Inches(1.22), Inches(6.6), Inches(8), Inches(0.4))
            _run(tfm.paragraphs[0], meta, font=BODY_FONT, size=12,
                 color=RGBColor(0xC6, 0xBF, 0xF2))
        return s

    # ── INDICE (light) ─────────────────────────────────────────────────────────
    def index(self, items, title="Indice"):
        s = self._new()
        _bg(s, PAPER)
        self._kicker(s, "Sommario", Inches(0.9), Inches(0.78), Inches(8))
        tb, tf = _txt(s, Inches(0.85), Inches(1.12), Inches(9), Inches(1.0))
        _run(tf.paragraphs[0], title, font=TITLE_FONT, size=40, color=INK, bold=True)
        # righe numerate su due colonne se > 4
        n = len(items)
        col_items = [items] if n <= 4 else [items[:(n + 1) // 2], items[(n + 1) // 2:]]
        col_x = [Inches(0.9), Inches(7.0)]
        start_y = 2.55
        row_h = 0.92
        for ci, col in enumerate(col_items):
            base = ci * 0  # same y start
            for i, label in enumerate(col):
                idx = (ci * len(col_items[0])) + i + 1 if n > 4 else i + 1
                y = Inches(start_y + i * row_h)
                # numero grande
                tbn, tfn = _txt(s, col_x[ci], y - Inches(0.06), Inches(1.0), Inches(0.9))
                _run(tfn.paragraphs[0], "%02d" % idx, font=TITLE_FONT, size=34,
                     color=PURPLE, bold=True)
                # etichetta
                tbl, tfl = _txt(s, col_x[ci] + Inches(1.0), y + Inches(0.05),
                                Inches(4.7), Inches(0.8), anchor=MSO_ANCHOR.TOP)
                _run(tfl.paragraphs[0], label, font=BODY_FONT, size=17, color=INK)
                # separatore sottile
                ln = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, col_x[ci],
                                        y + Inches(0.74), Inches(5.4), Pt(1))
                _solid(ln, RGBColor(0xEC, 0xEC, 0xF0))
        self._footer(s, dark=False)
        return s

    # ── DIVISORE DI SEZIONE (dark) ──────────────────────────────────────────────
    def section(self, number, title, subtitle=None):
        s = self._new()
        _bg(s, BLACK)
        # monogramma sfumato angolo basso dx
        _pic(s, _asset("mark_smoke.png"), Inches(9.4), Inches(2.2), h=Inches(6.6))
        # numero gigante arancione
        tbn, tfn = _txt(s, Inches(0.85), Inches(1.7), Inches(4), Inches(2.4))
        _run(tfn.paragraphs[0], "%02d" % number, font=TITLE_FONT, size=150,
             color=ORANGE, bold=True)
        # titolo
        tb, tf = _txt(s, Inches(0.9), Inches(4.5), Inches(9.6), Inches(2.0))
        p = tf.paragraphs[0]; p.line_spacing = 0.98
        _run(p, title, font=TITLE_FONT, size=46, color=WHITE, bold=True)
        if subtitle:
            tb2, tf2 = _txt(s, Inches(0.92), Inches(5.7), Inches(8.6), Inches(1.0))
            _run(tf2.paragraphs[0], subtitle, font=BODY_FONT, size=16,
                 color=RGBColor(0xC2, 0xC0, 0xCE))
        return s

    # ── CONTENUTO: titolo + testo + bullet (light) ──────────────────────────────
    def content(self, title, body=None, bullets=None, kicker=None):
        s = self._new()
        _bg(s, PAPER)
        self._kicker(s, kicker, Inches(0.9), Inches(0.78), Inches(10))
        tb, tf = _txt(s, Inches(0.85), Inches(1.12), Inches(11.4), Inches(1.1))
        p = tf.paragraphs[0]; p.line_spacing = 1.0
        _run(p, title, font=TITLE_FONT, size=34, color=INK, bold=True)
        bullets_y = 2.6
        if body:
            tbb, tfb = _txt(s, Inches(0.9), Inches(2.5), Inches(11.4), Inches(1.6))
            pb = tfb.paragraphs[0]; pb.line_spacing = 1.3
            _run(pb, body, font=BODY_FONT, size=16, color=GREY)
            # gap pulito sotto il paragrafo, in base alle righe stimate
            est_lines = max(1, 1 + len(body) // 95)
            bullets_y = 2.6 + 0.42 * est_lines + 0.45
        if bullets:
            tbl, tfl = _txt(s, Inches(0.9), Inches(min(bullets_y, 4.4)), Inches(11.0),
                            Inches(3.6))
            for i, b in enumerate(bullets):
                _bullet_para(tfl, b, first=(i == 0), size=16, color=INK)
        self._footer(s, dark=False)
        return s

    # ── CONTENUTO: due colonne (light) ──────────────────────────────────────────
    def two_column(self, title, left, right, kicker=None,
                   left_title=None, right_title=None):
        s = self._new()
        _bg(s, PAPER)
        self._kicker(s, kicker, Inches(0.9), Inches(0.78), Inches(10))
        tb, tf = _txt(s, Inches(0.85), Inches(1.12), Inches(11.4), Inches(1.1))
        _run(tf.paragraphs[0], title, font=TITLE_FONT, size=34, color=INK, bold=True)
        cols = [(Inches(0.9), left_title, left), (Inches(7.05), right_title, right)]
        for x, ctitle, content in cols:
            yy = 2.7
            if ctitle:
                tbt, tft = _txt(s, x, Inches(2.45), Inches(5.4), Inches(0.5))
                _run(tft.paragraphs[0], ctitle, font=BODY_FONT, size=18,
                     color=PURPLE, bold=True)
                yy = 3.05
            tbc, tfc = _txt(s, x, Inches(yy), Inches(5.4), Inches(3.4))
            if isinstance(content, (list, tuple)):
                for i, b in enumerate(content):
                    _bullet_para(tfc, b, first=(i == 0), size=15, color=INK)
            else:
                pc = tfc.paragraphs[0]; pc.line_spacing = 1.3
                _run(pc, content, font=BODY_FONT, size=15, color=GREY)
        self._footer(s, dark=False)
        return s

    # ── CONTENUTO: card (light) ──────────────────────────────────────────────────
    def cards(self, title, cards, kicker=None, intro=None, footnote=None):
        s = self._new()
        _bg(s, PAPER)
        self._kicker(s, kicker, Inches(0.9), Inches(0.78), Inches(10))
        tb, tf = _txt(s, Inches(0.85), Inches(1.12), Inches(11.4), Inches(1.1))
        _run(tf.paragraphs[0], title, font=TITLE_FONT, size=34, color=INK, bold=True)
        n = len(cards)
        gap = 0.4
        total_w = 11.53
        cw = (total_w - gap * (n - 1)) / n
        x0 = 0.9
        top = 2.75
        ch = 3.5
        if intro:
            tbi, tfi = _txt(s, Inches(0.9), Inches(2.4), Inches(11.4), Inches(0.7))
            pi = tfi.paragraphs[0]; pi.line_spacing = 1.25
            _run(pi, intro, font=BODY_FONT, size=14, color=GREY)
            top = 3.25
            ch = 3.0
        if footnote:
            ch = min(ch, 6.45 - top)
        for i, c in enumerate(cards):
            x = Inches(x0 + i * (cw + gap))
            card = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, Inches(top),
                                      Inches(cw), Inches(ch))
            _solid(card, MIST)
            card.adjustments[0] = 0.06
            card.shadow.inherit = False
            # numero/indice in pallino
            badge = s.shapes.add_shape(MSO_SHAPE.OVAL, x + Inches(0.4),
                                       Inches(top + 0.4), Inches(0.62), Inches(0.62))
            _solid(badge, PURPLE)
            btf = badge.text_frame
            btf.word_wrap = False
            bp = btf.paragraphs[0]; bp.alignment = PP_ALIGN.CENTER
            _run(bp, c.get("badge", str(i + 1)), font=TITLE_FONT, size=22,
                 color=WHITE, bold=True)
            # titolo card
            tbt, tft = _txt(s, x + Inches(0.4), Inches(top + 1.3),
                            Inches(cw - 0.8), Inches(0.7))
            _run(tft.paragraphs[0], c.get("title", ""), font=TITLE_FONT, size=21,
                 color=INK, bold=True)
            # testo card
            tbx, tfx = _txt(s, x + Inches(0.4), Inches(top + 2.0),
                            Inches(cw - 0.8), Inches(ch - 2.1))
            ptx = tfx.paragraphs[0]; ptx.line_spacing = 1.25
            _run(ptx, c.get("text", ""), font=BODY_FONT, size=14, color=GREY)
        if footnote:
            fy = top + ch + 0.25
            dot = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(0.9), Inches(fy + 0.04),
                                     Inches(0.12), Inches(0.12))
            _solid(dot, ORANGE)
            tbf, tff = _txt(s, Inches(1.16), Inches(fy), Inches(11.0), Inches(0.6))
            pf = tff.paragraphs[0]; pf.line_spacing = 1.2
            _run(pf, footnote, font=BODY_FONT, size=12, color=GREY)
        self._footer(s, dark=False)
        return s

    # ── CONTENUTO: statistiche / numeri (light) ─────────────────────────────────
    def stats(self, title, items, kicker=None, intro=None, footnote=None):
        s = self._new()
        _bg(s, PAPER)
        self._kicker(s, kicker, Inches(0.9), Inches(0.78), Inches(10))
        tb, tf = _txt(s, Inches(0.85), Inches(1.12), Inches(11.4), Inches(1.1))
        _run(tf.paragraphs[0], title, font=TITLE_FONT, size=34, color=INK, bold=True)
        has_desc = any(len(it) > 2 for it in items)
        big_y = 3.0
        if intro:
            tbi, tfi = _txt(s, Inches(0.9), Inches(2.35), Inches(11.4), Inches(0.7))
            pi = tfi.paragraphs[0]; pi.line_spacing = 1.25
            _run(pi, intro, font=BODY_FONT, size=14, color=GREY)
            big_y = 3.25
        n = len(items)
        cw = 11.53 / n
        big_sz = 56 if has_desc else 68
        for i, it in enumerate(items):
            big, label = it[0], it[1]
            desc = it[2] if len(it) > 2 else None
            x = Inches(0.9 + i * cw)
            tbb, tfb = _txt(s, x, Inches(big_y), Inches(cw - 0.3), Inches(1.1))
            _run(tfb.paragraphs[0], big, font=TITLE_FONT, size=big_sz,
                 color=PURPLE if i % 2 == 0 else ORANGE, bold=True)
            tbl, tfl = _txt(s, x, Inches(big_y + 1.12), Inches(cw - 0.45), Inches(0.5))
            _run(tfl.paragraphs[0], label, font=BODY_FONT, size=15, color=INK, bold=True)
            if desc:
                tbd, tfd = _txt(s, x, Inches(big_y + 1.6), Inches(cw - 0.45), Inches(1.4))
                pd = tfd.paragraphs[0]; pd.line_spacing = 1.2
                _run(pd, desc, font=BODY_FONT, size=12.5, color=GREY)
        if footnote:
            dot = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(0.9), Inches(6.36),
                                     Inches(0.12), Inches(0.12))
            _solid(dot, ORANGE)
            tbf, tff = _txt(s, Inches(1.16), Inches(6.3), Inches(11.0), Inches(0.5))
            _run(tff.paragraphs[0], footnote, font=BODY_FONT, size=12.5,
                 color=INK, bold=True)
        self._footer(s, dark=False)
        return s

    def _head(self, s, title, kicker):
        _bg(s, PAPER)
        self._kicker(s, kicker, Inches(0.9), Inches(0.78), Inches(11))
        tb, tf = _txt(s, Inches(0.85), Inches(1.12), Inches(11.6), Inches(1.0))
        _run(tf.paragraphs[0], title, font=TITLE_FONT, size=32, color=INK, bold=True)

    # ── SWATCH a gruppi (es. status con colore) ─────────────────────────────────
    def swatch_groups(self, title, groups, kicker=None, intro=None):
        s = self._new(); self._head(s, title, kicker)
        top = 2.35
        if intro:
            tbi, tfi = _txt(s, Inches(0.9), Inches(2.2), Inches(11.4), Inches(0.7))
            pi = tfi.paragraphs[0]; pi.line_spacing = 1.25
            _run(pi, intro, font=BODY_FONT, size=14, color=GREY)
            top = 3.0
        n = len(groups); colw = 11.53 / n; x0 = 0.9
        for gi, g in enumerate(groups):
            x = x0 + gi * colw
            tbg, tfg = _txt(s, Inches(x), Inches(top), Inches(colw - 0.3), Inches(0.45))
            _run(tfg.paragraphs[0], g["label"].upper(), font=TITLE_FONT, size=12,
                 color=PURPLE, bold=True, spacing=1.5)
            y = top + 0.62
            for it in g["items"]:
                chip = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x),
                                          Inches(y), Inches(0.32), Inches(0.32))
                _solid(chip, _hex(it["color"])); chip.adjustments[0] = 0.35
                chip.shadow.inherit = False
                tbn, tfn = _txt(s, Inches(x + 0.48), Inches(y - 0.04),
                                Inches(colw - 0.6), Inches(0.34))
                pn = tfn.paragraphs[0]
                _run(pn, it["name"], font=BODY_FONT, size=14, color=INK, bold=True)
                if it.get("code"):
                    _run(pn, "   " + it["code"], font=BODY_FONT, size=10.5, color=FAINT)
                if it.get("desc"):
                    tbd, tfd = _txt(s, Inches(x + 0.48), Inches(y + 0.28),
                                    Inches(colw - 0.7), Inches(0.5))
                    pd = tfd.paragraphs[0]; pd.line_spacing = 1.1
                    _run(pd, it["desc"], font=BODY_FONT, size=11, color=GREY)
                y += 0.82
        self._footer(s)
        return s

    # ── GRIGLIA di swatch circolari (es. scala cluster) ─────────────────────────
    def swatch_grid(self, title, items, kicker=None, cols=3, intro=None):
        s = self._new(); self._head(s, title, kicker)
        top = 2.55
        if intro:
            tbi, tfi = _txt(s, Inches(0.9), Inches(2.2), Inches(11.4), Inches(0.7))
            pi = tfi.paragraphs[0]; pi.line_spacing = 1.25
            _run(pi, intro, font=BODY_FONT, size=14, color=GREY)
            top = 3.05
        cw = 11.53 / cols
        rh = 1.55
        for i, it in enumerate(items):
            r = i // cols; c = i % cols
            x = 0.9 + c * cw; y = top + r * rh
            d = it.get("d", 0.7)
            circ = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x), Inches(y),
                                      Inches(d), Inches(d))
            _solid(circ, _hex(it["color"])); circ.shadow.inherit = False
            tx = x + d + 0.25
            tbn, tfn = _txt(s, Inches(tx), Inches(y - 0.04), Inches(cw - d - 0.45),
                            Inches(0.34))
            _run(tfn.paragraphs[0], it["name"], font=BODY_FONT, size=14,
                 color=INK, bold=True)
            if it.get("sub"):
                tbs, tfs = _txt(s, Inches(tx), Inches(y + 0.28),
                                Inches(cw - d - 0.45), Inches(0.3))
                _run(tfs.paragraphs[0], it["sub"], font=BODY_FONT, size=11.5, color=GREY)
            if it.get("code"):
                tbc, tfc = _txt(s, Inches(tx), Inches(y + 0.55),
                                Inches(cw - d - 0.45), Inches(0.3))
                _run(tfc.paragraphs[0], it["code"], font=BODY_FONT, size=10.5, color=FAINT)
        self._footer(s)
        return s

    # ── TAG LIST (es. tipologie) ─────────────────────────────────────────────────
    def taglist(self, title, groups, kicker=None):
        s = self._new(); self._head(s, title, kicker)
        n = len(groups); colw = 11.53 / n; x0 = 0.9
        for gi, g in enumerate(groups):
            bx = x0 + gi * colw
            tbg, tfg = _txt(s, Inches(bx), Inches(2.4), Inches(colw - 0.4), Inches(0.4))
            _run(tfg.paragraphs[0], g["label"].upper(), font=TITLE_FONT, size=12,
                 color=PURPLE, bold=True, spacing=1.5)
            # flusso di pill
            px, py = bx, 3.0
            line_h = 0.5
            maxx = bx + colw - 0.4
            for name in g["items"]:
                w = 0.22 + 0.092 * len(name)
                if px + w > maxx and px > bx:
                    px = bx; py += line_h
                pill = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(px),
                                          Inches(py), Inches(w), Inches(0.38))
                _solid(pill, MIST); pill.adjustments[0] = 0.5; pill.shadow.inherit = False
                ptf = pill.text_frame
                for m in ("margin_left", "margin_right", "margin_top", "margin_bottom"):
                    setattr(ptf, m, 0)
                pp = ptf.paragraphs[0]; pp.alignment = PP_ALIGN.CENTER
                _run(pp, name, font=BODY_FONT, size=12, color=INK)
                px += w + 0.12
        self._footer(s)
        return s

    # ── TABELLA di confronto ─────────────────────────────────────────────────────
    def table(self, title, headers, rows, kicker=None, col_widths=None, intro=None):
        s = self._new(); self._head(s, title, kicker)
        ncol = len(headers); nrow = len(rows) + 1
        top_in = 2.4
        if intro:
            tbi, tfi = _txt(s, Inches(0.9), Inches(2.02), Inches(11.5), Inches(0.5))
            pi = tfi.paragraphs[0]; pi.line_spacing = 1.2
            _run(pi, intro, font=BODY_FONT, size=13, color=GREY)
            top_in = 2.55
        rh = min(0.5, (6.55 - top_in) / nrow)
        left, top, total_w = Inches(0.85), Inches(top_in), Inches(11.6)
        height = Inches(rh * nrow)
        gtbl = s.shapes.add_table(nrow, ncol, left, top, total_w, height).table
        gtbl.first_row = False; gtbl.horz_banding = False
        if col_widths:
            for i, w in enumerate(col_widths):
                gtbl.columns[i].width = Inches(w)
        # header
        for c, h in enumerate(headers):
            cell = gtbl.cell(0, c)
            cell.fill.solid(); cell.fill.fore_color.rgb = PURPLE
            cell.margin_left = Inches(0.12); cell.margin_right = Inches(0.1)
            cell.margin_top = Inches(0.06); cell.margin_bottom = Inches(0.06)
            p = cell.text_frame.paragraphs[0]
            _run(p, h, font=TITLE_FONT, size=13, color=WHITE, bold=True)
        # body
        for r, row in enumerate(rows, start=1):
            for c, val in enumerate(row):
                cell = gtbl.cell(r, c)
                cell.fill.solid()
                cell.fill.fore_color.rgb = WHITE if r % 2 else MIST
                cell.margin_left = Inches(0.12); cell.margin_right = Inches(0.1)
                cell.margin_top = Inches(0.05); cell.margin_bottom = Inches(0.05)
                p = cell.text_frame.paragraphs[0]; p.line_spacing = 1.05
                lead = (c == 0)
                _run(p, val, font=BODY_FONT, size=11,
                     color=INK if lead else GREY, bold=lead)
        self._footer(s)
        return s

    # ── CHIUSURA (dark) ──────────────────────────────────────────────────────────
    def closing(self, big="buon lavoro", small=None):
        s = self._new()
        _bg(s, BLACK)
        _pic(s, _asset("mark_smoke.png"), Inches(8.9), Inches(-0.9), h=Inches(9.4))
        w = Inches(1.7); h = Emu(int(w * 120 / 357))
        _pic(s, _asset("type_white.png"), Inches(0.85), Inches(0.85), w=w, h=h)
        tb, tf = _txt(s, Inches(0.85), Inches(2.9), Inches(10.5), Inches(2.2),
                      anchor=MSO_ANCHOR.MIDDLE)
        p = tf.paragraphs[0]; p.line_spacing = 0.98
        _run(p, big, font=TITLE_FONT, size=66, color=WHITE, bold=True)
        dot = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(0.92), Inches(5.2),
                                 Inches(0.16), Inches(0.16))
        _solid(dot, ORANGE)
        if small:
            tbs, tfs = _txt(s, Inches(1.22), Inches(5.14), Inches(8), Inches(0.4))
            _run(tfs.paragraphs[0], small, font=BODY_FONT, size=13,
                 color=RGBColor(0x9A, 0x99, 0xA6))
        return s

    # ── SALVATAGGIO + EMBED FONT ─────────────────────────────────────────────────
    def save(self, path, embed=True):
        self.prs.save(path)
        if embed:
            try:
                embed_fonts(path)
            except Exception as e:           # non bloccare il salvataggio
                print("[warn] embed font non riuscito:", e)
        return path


# ─────────────────────────────────────────────────────────────────────────────
#  EMBEDDING FONT NEL PACCHETTO PPTX
# ─────────────────────────────────────────────────────────────────────────────
def embed_fonts(pptx_path, fonts_dir=FONTS):
    """Incorpora Oswald e Inter (Regular+Bold) nel file .pptx."""
    families = [
        ("Oswald", {"regular": "Oswald-Regular.ttf", "bold": "Oswald-Bold.ttf"}),
        ("Inter",  {"regular": "Inter-Regular.ttf",  "bold": "Inter-Bold.ttf"}),
    ]
    # verifica presenza file
    for _, slots in families:
        for fn in slots.values():
            if not os.path.exists(os.path.join(fonts_dir, fn)):
                raise FileNotFoundError(os.path.join(fonts_dir, fn))

    tmp = pptx_path + ".tmp"
    with zipfile.ZipFile(pptx_path, "r") as zin:
        names = zin.namelist()
        data = {n: zin.read(n) for n in names}

    # 1) font binari
    font_index = 1
    rel_entries = []        # (rId, target)
    embed_xml = []          # blocchi <p:embeddedFont>
    rid_base = 9000
    for typeface, slots in families:
        block = ['<p:embeddedFont><p:font typeface="%s"/>' % typeface]
        for style in ("regular", "bold"):
            fn = slots[style]
            target = "fonts/font%d.fntdata" % font_index
            data["ppt/" + target] = open(os.path.join(fonts_dir, fn), "rb").read()
            rid = "rId%d" % (rid_base + font_index)
            rel_entries.append((rid, target))
            block.append('<p:%s r:id="%s"/>' % (style, rid))
            font_index += 1
        block.append("</p:embeddedFont>")
        embed_xml.append("".join(block))

    # 2) [Content_Types].xml  → default per fntdata
    ct = data["[Content_Types].xml"].decode("utf-8")
    if "fntdata" not in ct:
        ins = '<Default Extension="fntdata" ContentType="application/x-fontdata"/>'
        ct = ct.replace("</Types>", ins + "</Types>")
    data["[Content_Types].xml"] = ct.encode("utf-8")

    # 3) presentation.xml.rels → relazioni font
    rels_path = "ppt/_rels/presentation.xml.rels"
    rels = data[rels_path].decode("utf-8")
    import re as _re
    rels = _re.sub(r'<Relationship[^>]*relationships/font"[^>]*/>', "", rels)
    rel_ins = "".join(
        '<Relationship Id="%s" '
        'Type="http://schemas.openxmlformats.org/officeDocument/2006/relationships/font" '
        'Target="%s"/>' % (rid, tgt) for rid, tgt in rel_entries
    )
    rels = rels.replace("</Relationships>", rel_ins + "</Relationships>")
    data[rels_path] = rels.encode("utf-8")

    # 4) presentation.xml → attributi + <p:embeddedFontLst>
    pres_path = "ppt/presentation.xml"
    pres = data[pres_path].decode("utf-8")
    # attributi sul tag <p:presentation ...> (idempotente)
    import re
    m = re.search(r"<p:presentation\b[^>]*>", pres)
    open_tag = m.group(0)
    new_tag = open_tag
    for attr, val in (("embedTrueTypeFonts", "1"), ("saveSubsetFonts", "0")):
        if (attr + "=") not in new_tag:
            new_tag = new_tag[:-1] + ' %s="%s"' % (attr, val) + ">"
    pres = pres[:m.start()] + new_tag + pres[m.end():]
    # evita duplicati della lista font se gia' presente
    pres = re.sub(r"<p:embeddedFontLst>.*?</p:embeddedFontLst>", "", pres, flags=re.S)
    lst = "<p:embeddedFontLst>" + "".join(embed_xml) + "</p:embeddedFontLst>"
    # inserire DOPO </p:notesSz> (ordine schema CT_Presentation)
    if "</p:notesSz>" in pres:
        pres = pres.replace("</p:notesSz>", "</p:notesSz>" + lst, 1)
    elif "<p:notesSz" in pres:
        pres = re.sub(r"(<p:notesSz[^>]*/>)", r"\1" + lst, pres, count=1)
    else:  # fallback: dopo sldIdLst
        pres = pres.replace("</p:sldIdLst>", "</p:sldIdLst>" + lst, 1)
    data[pres_path] = pres.encode("utf-8")

    # riscrivi lo zip
    with zipfile.ZipFile(tmp, "w", zipfile.ZIP_DEFLATED) as zout:
        for n, b in data.items():
            zout.writestr(n, b)
    shutil.move(tmp, pptx_path)
    return pptx_path


# ─────────────────────────────────────────────────────────────────────────────
#  GENERAZIONE DA SPEC (dict / JSON) — per impaginazione automatica
# ─────────────────────────────────────────────────────────────────────────────
def build_from_spec(spec, out_path):
    """
    spec = {
      "slides": [
        {"type": "cover", "title": "...", "subtitle": "...", "meta": "...",
         "kicker": "..."},
        {"type": "index", "items": ["A", "B", "C"]},
        {"type": "section", "number": 1, "title": "...", "subtitle": "..."},
        {"type": "content", "title": "...", "body": "...",
         "bullets": ["..."], "kicker": "..."},
        {"type": "two_column", "title": "...", "left": [...], "right": [...],
         "left_title": "...", "right_title": "...", "kicker": "..."},
        {"type": "cards", "title": "...", "kicker": "...",
         "cards": [{"title": "...", "text": "...", "badge": "1"}]},
        {"type": "stats", "title": "...", "kicker": "...",
         "items": [["120", "etichetta"]]},
        {"type": "closing", "big": "buon lavoro", "small": "..."}
      ]
    }
    """
    d = GravityDeck()
    for sl in spec.get("slides", []):
        t = sl.get("type")
        if t == "cover":
            d.cover(sl["title"], sl.get("subtitle"), sl.get("meta"), sl.get("kicker"))
        elif t == "index":
            d.index(sl["items"], sl.get("title", "Indice"))
        elif t == "section":
            d.section(sl["number"], sl["title"], sl.get("subtitle"))
        elif t == "content":
            d.content(sl["title"], sl.get("body"), sl.get("bullets"), sl.get("kicker"))
        elif t == "two_column":
            d.two_column(sl["title"], sl["left"], sl["right"], sl.get("kicker"),
                         sl.get("left_title"), sl.get("right_title"))
        elif t == "cards":
            d.cards(sl["title"], sl["cards"], sl.get("kicker"))
        elif t == "stats":
            d.stats(sl["title"], sl["items"], sl.get("kicker"))
        elif t == "closing":
            d.closing(sl.get("big", "buon lavoro"), sl.get("small"))
        else:
            raise ValueError("tipo slide sconosciuto: %r" % t)
    return d.save(out_path)


if __name__ == "__main__":
    import json
    import sys
    if len(sys.argv) == 3:
        spec = json.load(open(sys.argv[1], encoding="utf-8"))
        build_from_spec(spec, sys.argv[2])
        print("creato", sys.argv[2])
    else:
        print("uso: python gravity_deck.py spec.json output.pptx")

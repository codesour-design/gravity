#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
GRAVITY — costruzione del TEMA con MASTRO + LAYOUT (segnaposto).

A differenza di gravity_deck.py (che disegna sulle slide), qui la grafica vive
nello slide master e nei layout: cover, sezione, contenuto, due colonne,
chiusura. I colori brand finiscono nel tema, i font major/minor sono Oswald e
Inter. Il file, caricato su Google Slides convertito, espone i layout
nell'editor mastro e nel selettore "Applica layout".
"""
import os
import re
import shutil
import zipfile
from lxml import etree
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor
from pptx.oxml.ns import qn
from pptx.shapes.shapetree import SlideShapes

import gravity_deck as G  # riusa colori, embed_fonts


def _shp(layout):
    """Wrapper che espone add_shape/add_picture/add_textbox su un layout."""
    return SlideShapes(layout.shapes._spTree, layout)

HERE = os.path.dirname(os.path.abspath(__file__))
ASSETS = os.path.join(HERE, "assets")
A = lambda n: os.path.join(ASSETS, n)

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)
EMU = 914400
NS_A = "http://schemas.openxmlformats.org/drawingml/2006/main"


def hexstr(c):
    return "%02X%02X%02X" % (c[0], c[1], c[2])


# ── helper grafica su layout ────────────────────────────────────────────────
def _send_back(shapes, shp):
    el = shp._element
    el.getparent().remove(el)
    shapes._spTree.insert(2, el)


def _bg(layout, color):
    sh = _shp(layout)
    r = sh.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, SLIDE_H)
    r.fill.solid(); r.fill.fore_color.rgb = color
    r.line.fill.background(); r.shadow.inherit = False
    _send_back(sh, r)
    return r


def _pic(layout, path, x, y, w=None, h=None):
    return _shp(layout).add_picture(path, x, y, width=w, height=h)


def _static_text(layout, text, x, y, w, h, font, size, color, bold=False,
                 align=PP_ALIGN.LEFT, spacing=None, line=None):
    tb = _shp(layout).add_textbox(x, y, w, h)
    tf = tb.text_frame; tf.word_wrap = True
    for m in ("margin_left", "margin_right", "margin_top", "margin_bottom"):
        setattr(tf, m, 0)
    p = tf.paragraphs[0]; p.alignment = align
    if line:
        p.line_spacing = line
    r = p.add_run(); r.text = text
    r.font.name = font; r.font.size = Pt(size); r.font.bold = bold
    r.font.color.rgb = color
    if spacing is not None:
        r._r.get_or_add_rPr().set("spc", str(int(spacing * 100)))
    return tb


def _dot(layout, x, y, d, color):
    o = _shp(layout).add_shape(MSO_SHAPE.OVAL, x, y, Inches(d), Inches(d))
    o.fill.solid(); o.fill.fore_color.rgb = color; o.line.fill.background()
    o.shadow.inherit = False
    return o


def place(ph, x, y, w, h):
    ph.left, ph.top, ph.width, ph.height = x, y, w, h


def style_ph(ph, *, size=None, bold=None, color=None, font=None,
             align=None, bullets=False, line=None, anchor=None):
    """Imposta lo stile di default del segnaposto a livello di LAYOUT."""
    tf = ph.text_frame
    if anchor is not None:
        tf.vertical_anchor = anchor
    for m in ("margin_left", "margin_right", "margin_top", "margin_bottom"):
        setattr(tf, m, 0)
    txBody = tf._txBody
    bodyPr = txBody.find(qn("a:bodyPr"))
    lstStyle = txBody.find(qn("a:lstStyle"))
    if lstStyle is None:
        lstStyle = txBody.makeelement(qn("a:lstStyle"), {})
        bodyPr.addnext(lstStyle)
    for e in lstStyle.findall(qn("a:lvl1pPr")):
        lstStyle.remove(e)
    lvl = lstStyle.makeelement(qn("a:lvl1pPr"), {})
    if align:
        lvl.set("algn", align)
    if bullets:
        lvl.set("marL", str(int(0.30 * EMU)))
        lvl.set("indent", str(int(-0.30 * EMU)))
    else:
        lvl.set("marL", "0"); lvl.set("indent", "0")
    # ordine schema: lnSpc, (spc), buClr, buSz, buFont, bu*, defRPr
    if line:
        ln = lvl.makeelement(qn("a:lnSpc"), {})
        ln.append(lvl.makeelement(qn("a:spcPct"), {"val": str(int(line * 100000))}))
        lvl.append(ln)
    if bullets:
        buClr = lvl.makeelement(qn("a:buClr"), {})
        buClr.append(lvl.makeelement(qn("a:srgbClr"), {"val": hexstr(G.ORANGE)}))
        lvl.append(buClr)
        lvl.append(lvl.makeelement(qn("a:buSzPct"), {"val": "70000"}))
        lvl.append(lvl.makeelement(qn("a:buFont"), {"typeface": "Arial"}))
        lvl.append(lvl.makeelement(qn("a:buChar"), {"char": "▪"}))
    else:
        lvl.append(lvl.makeelement(qn("a:buNone"), {}))
    defRPr = lvl.makeelement(qn("a:defRPr"), {})
    if size:
        defRPr.set("sz", str(int(size * 100)))
    if bold is not None:
        defRPr.set("b", "1" if bold else "0")
    if color:
        sf = defRPr.makeelement(qn("a:solidFill"), {})
        sf.append(defRPr.makeelement(qn("a:srgbClr"), {"val": hexstr(color)}))
        defRPr.append(sf)
    if font:
        defRPr.append(defRPr.makeelement(qn("a:latin"), {"typeface": font}))
    lvl.append(defRPr)
    lstStyle.append(lvl)


def set_name(layout, name):
    layout._element.cSld.set("name", name)


def fill(slide, idx, text):
    ph = slide.placeholders[idx]
    ph.text = text
    return ph


# ── costruzione ─────────────────────────────────────────────────────────────
def build(path, embed=True):
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    master = prs.slide_masters[0]
    L = master.slide_layouts

    PURPLE, ORANGE, BLACK, INK, GREY = G.PURPLE, G.ORANGE, G.BLACK, G.INK, G.GREY
    WHITE = G.WHITE
    LAV = RGBColor(0xDA, 0xD3, 0xFF)
    META = RGBColor(0xC6, 0xBF, 0xF2)
    wm_w = Inches(1.7); wm_h = Emu(int(wm_w * 120 / 357))
    foot_w = Inches(1.05); foot_h = Emu(int(foot_w * 120 / 357))

    # idx default: 0 Title, 1 Title+Content, 2 Section, 3 Two Content, 5 Title Only
    # ===== LAYOUT 0 — COVER (purple) =====
    lo = L[0]; set_name(lo, "Gravity – Cover")
    _bg(lo, PURPLE)
    _pic(lo, A("mark_deep.png"), Inches(8.7), Inches(-1.1), h=Inches(9.6))
    _pic(lo, A("type_white.png"), Inches(0.85), Inches(0.8), w=wm_w, h=wm_h)
    _static_text(lo, "TIPO DI DOCUMENTO", Inches(0.9), Inches(2.55), Inches(8),
                 Inches(0.34), G.TITLE_FONT, 13, ORANGE, bold=True, spacing=2.5)
    _dot(lo, Inches(0.92), Inches(6.66), 0.16, ORANGE)
    _static_text(lo, "CodeSour · Gravity · 2026", Inches(1.22), Inches(6.6),
                 Inches(8), Inches(0.4), G.BODY_FONT, 12, META)
    t = lo.placeholders[0]   # ctrTitle
    place(t, Inches(0.85), Inches(2.85), Inches(9.4), Inches(2.4))
    style_ph(t, size=54, bold=True, color=WHITE, font=G.TITLE_FONT,
             align="l", line=0.98, anchor=MSO_ANCHOR.TOP)
    st = lo.placeholders[1]  # subTitle
    place(st, Inches(0.9), Inches(5.35), Inches(8.6), Inches(1.1))
    style_ph(st, size=18, bold=False, color=LAV, font=G.BODY_FONT,
             align="l", line=1.25)

    # ===== LAYOUT 2 — SEZIONE (dark) =====
    lo = L[2]; set_name(lo, "Gravity – Sezione")
    _bg(lo, BLACK)
    _pic(lo, A("mark_smoke.png"), Inches(9.4), Inches(2.2), h=Inches(6.6))
    _static_text(lo, "01", Inches(0.85), Inches(1.55), Inches(4), Inches(2.4),
                 G.TITLE_FONT, 150, ORANGE, bold=True)
    ti = lo.placeholders[0]   # title
    place(ti, Inches(0.9), Inches(4.4), Inches(9.6), Inches(1.3))
    style_ph(ti, size=44, bold=True, color=WHITE, font=G.TITLE_FONT,
             align="l", line=0.98)
    tx = lo.placeholders[1]   # text
    place(tx, Inches(0.92), Inches(5.65), Inches(8.6), Inches(1.0))
    style_ph(tx, size=16, bold=False, color=RGBColor(0xC2, 0xC0, 0xCE),
             font=G.BODY_FONT, align="l", line=1.2)

    # ===== LAYOUT 1 — CONTENUTO (light) =====
    lo = L[1]; set_name(lo, "Gravity – Contenuto")
    _bg(lo, WHITE)
    _pic(lo, A("type_purple.png"), Inches(0.6), Inches(6.88), w=foot_w, h=foot_h)
    _static_text(lo, "SEZIONE", Inches(0.9), Inches(0.78), Inches(10),
                 Inches(0.34), G.TITLE_FONT, 13, ORANGE, bold=True, spacing=2.5)
    ti = lo.placeholders[0]
    place(ti, Inches(0.85), Inches(1.12), Inches(11.5), Inches(1.0))
    style_ph(ti, size=34, bold=True, color=INK, font=G.TITLE_FONT, align="l",
             line=1.0)
    bd = lo.placeholders[1]
    place(bd, Inches(0.9), Inches(2.55), Inches(11.4), Inches(4.0))
    style_ph(bd, size=16, bold=False, color=INK, font=G.BODY_FONT, align="l",
             bullets=True, line=1.2)

    # ===== LAYOUT 3 — DUE COLONNE (light) =====
    lo = L[3]; set_name(lo, "Gravity – Due colonne")
    _bg(lo, WHITE)
    _pic(lo, A("type_purple.png"), Inches(0.6), Inches(6.88), w=foot_w, h=foot_h)
    _static_text(lo, "SEZIONE", Inches(0.9), Inches(0.78), Inches(10),
                 Inches(0.34), G.TITLE_FONT, 13, ORANGE, bold=True, spacing=2.5)
    ti = lo.placeholders[0]
    place(ti, Inches(0.85), Inches(1.12), Inches(11.5), Inches(1.0))
    style_ph(ti, size=34, bold=True, color=INK, font=G.TITLE_FONT, align="l",
             line=1.0)
    # due body placeholders (idx 1 e 2 nel layout Two Content)
    phs = [p for p in lo.placeholders if p.placeholder_format.idx in (1, 2)]
    geos = [(Inches(0.9), Inches(2.55), Inches(5.4), Inches(3.9)),
            (Inches(7.05), Inches(2.55), Inches(5.4), Inches(3.9))]
    for ph, g in zip(phs, geos):
        place(ph, *g)
        style_ph(ph, size=15, bold=False, color=INK, font=G.BODY_FONT,
                 align="l", bullets=True, line=1.2)

    # ===== LAYOUT 5 — CHIUSURA (dark) =====
    lo = L[5]; set_name(lo, "Gravity – Chiusura")
    _bg(lo, BLACK)
    _pic(lo, A("mark_smoke.png"), Inches(8.9), Inches(-0.9), h=Inches(9.4))
    _pic(lo, A("type_white.png"), Inches(0.85), Inches(0.85), w=wm_w, h=wm_h)
    _dot(lo, Inches(0.92), Inches(5.2), 0.16, ORANGE)
    _static_text(lo, "Gravity · CodeSour", Inches(1.22), Inches(5.14),
                 Inches(8), Inches(0.4), G.BODY_FONT, 13, RGBColor(0x9A, 0x99, 0xA6))
    ti = lo.placeholders[0]
    place(ti, Inches(0.85), Inches(2.9), Inches(10.5), Inches(1.8))
    style_ph(ti, size=62, bold=True, color=WHITE, font=G.TITLE_FONT, align="l",
             line=0.98, anchor=MSO_ANCHOR.MIDDLE)

    keep = {"Gravity – Cover", "Gravity – Sezione", "Gravity – Contenuto",
            "Gravity – Due colonne", "Gravity – Chiusura"}

    # ── slide demo (una per layout) ──
    s = prs.slides.add_slide(L[0]); fill(s, 0, "Titolo del documento")
    fill(s, 1, "Sottotitolo o claim su una o due righe.")
    s = prs.slides.add_slide(L[2]); fill(s, 0, "Contesto e obiettivi")
    fill(s, 1, "Il punto di partenza e cosa vogliamo ottenere.")
    s = prs.slides.add_slide(L[1]); fill(s, 0, "Il problema da risolvere")
    bd = s.placeholders[1]; tf = bd.text_frame
    tf.text = "Primo punto chiave, conciso e azionabile."
    for txt in ["Secondo punto che aggiunge una sfumatura.",
                "Terzo punto a chiusura del ragionamento."]:
        tf.add_paragraph().text = txt
    s = prs.slides.add_slide(L[3]); fill(s, 0, "Due binari complementari")
    phs = [p for p in s.placeholders if p.placeholder_format.idx in (1, 2)]
    phs[0].text_frame.text = "Sul prototipo: traccia il flusso."
    phs[0].text_frame.add_paragraph().text = "Inspector dei componenti."
    phs[1].text_frame.text = "Su Jira: lavoro in ticket."
    phs[1].text_frame.add_paragraph().text = "Stimabile e tracciabile."
    s = prs.slides.add_slide(L[5]); fill(s, 0, "buon lavoro")

    prs.save(path)
    # post-process: tema (colori+font), rimozione layout inutili, embed font
    _post(path, keep, embed=embed)
    return path


# ── post-process zip: tema + rimozione layout + embed font ──────────────────
def _post(path, keep_names, embed=True):
    with zipfile.ZipFile(path) as z:
        data = {n: z.read(n) for n in z.namelist()}

    # 1) THEME: colori brand + font major/minor
    tp = "ppt/theme/theme1.xml"
    if tp in data:
        root = etree.fromstring(data[tp])
        a = "{%s}" % NS_A
        clr = root.find(".//%sclrScheme" % a)
        def setclr(slot, hexv):
            el = clr.find("%s%s" % (a, slot))
            if el is None:
                return
            for ch in list(el):
                el.remove(ch)
            srgb = el.makeelement("%ssrgbClr" % a, {"val": hexv})
            el.append(srgb)
        setclr("accent1", hexstr(G.PURPLE))
        setclr("accent2", hexstr(G.ORANGE))
        setclr("dk2", "0A0A0A")
        setclr("lt2", "F5F5F5")
        fs = root.find(".//%sfontScheme" % a)
        for which, face in (("majorFont", "Oswald"), ("minorFont", "Inter")):
            latin = fs.find("%s%s/%slatin" % (a, which, a))
            if latin is not None:
                latin.set("typeface", face)
        data[tp] = etree.tostring(root, xml_declaration=True, encoding="UTF-8",
                                  standalone=True)

    # 2) rimuovi i layout non-Gravity
    mp = "ppt/slideMasters/slideMaster1.xml"
    mrp = "ppt/slideMasters/_rels/slideMaster1.xml.rels"
    P = "{http://schemas.openxmlformats.org/presentationml/2006/main}"
    R = "{http://schemas.openxmlformats.org/officeDocument/2006/relationships}"
    mroot = etree.fromstring(data[mp])
    rroot = etree.fromstring(data[mrp])
    rid_to_target = {r.get("Id"): r.get("Target") for r in rroot
                     if r.get("Type", "").endswith("/slideLayout")}
    # nome di ogni layout
    def layout_name(target):
        lp = "ppt/" + target.replace("../", "")
        try:
            lr = etree.fromstring(data[lp])
            return lr.find("%scSld" % P).get("name")
        except Exception:
            return None
    remove_rids = []
    for lst in mroot.findall("%ssldLayoutIdLst/%ssldLayoutId" % (P, P)):
        rid = lst.get("%sid" % R)
        tgt = rid_to_target.get(rid)
        nm = layout_name(tgt) if tgt else None
        if nm not in keep_names:
            remove_rids.append((rid, lst, tgt))
    for rid, lst, tgt in remove_rids:
        lst.getparent().remove(lst)                      # da sldLayoutIdLst
        for r in list(rroot):
            if r.get("Id") == rid:
                rroot.remove(r)
        # rimuovi parti del layout
        lp = "ppt/" + tgt.replace("../", "")
        data.pop(lp, None)
        lrels = lp.replace("slideLayouts/", "slideLayouts/_rels/") + ".rels"
        data.pop(lrels, None)
    data[mp] = etree.tostring(mroot, xml_declaration=True, encoding="UTF-8",
                              standalone=True)
    data[mrp] = etree.tostring(rroot, xml_declaration=True, encoding="UTF-8",
                               standalone=True)
    # pulizia [Content_Types].xml dalle override dei layout rimossi
    ct = data["[Content_Types].xml"].decode("utf-8")
    for rid, lst, tgt in remove_rids:
        part = "/ppt/" + tgt.replace("../", "")
        ct = re.sub(r'<Override PartName="%s"[^>]*/>' % re.escape(part), "", ct)
    data["[Content_Types].xml"] = ct.encode("utf-8")

    # riscrivi
    tmp = path + ".tmp"
    with zipfile.ZipFile(tmp, "w", zipfile.ZIP_DEFLATED) as zout:
        for n, b in data.items():
            zout.writestr(n, b)
    shutil.move(tmp, path)

    # 3) embed font (per uso PPTX; su Slides Oswald/Inter sono Google Font)
    if embed:
        try:
            G.embed_fonts(path)
        except Exception as e:
            print("[warn] embed:", e)
    return path


if __name__ == "__main__":
    import sys
    out = sys.argv[1] if len(sys.argv) > 1 else "Gravity_Master.pptx"
    build(out)
    print("creato", out)
